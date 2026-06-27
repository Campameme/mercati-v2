"""
Genera presentazione.pptx editabile a partire dal contenuto delle slide HTML.
Stile coerente con IMercati: palette cream/olive/ink, Fraunces serif + Inter sans.

Esecuzione:
    python build_pptx.py
Output:
    presentazione.pptx (nella stessa cartella)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# Palette IMercati
CREAM_50 = RGBColor(0xFA, 0xF7, 0xF1)
CREAM_100 = RGBColor(0xF4, 0xED, 0xE0)
CREAM_200 = RGBColor(0xEB, 0xE0, 0xCC)
CREAM_300 = RGBColor(0xDD, 0xCD, 0xB0)
INK = RGBColor(0x2A, 0x26, 0x20)
INK_SOFT = RGBColor(0x4A, 0x43, 0x39)
INK_MUTED = RGBColor(0x7A, 0x6F, 0x60)
OLIVE_50 = RGBColor(0xF0, 0xF1, 0xE4)
OLIVE_300 = RGBColor(0xB8, 0xC1, 0x82)
OLIVE_500 = RGBColor(0x7D, 0x8F, 0x4E)
OLIVE_600 = RGBColor(0x5D, 0x6E, 0x3B)
OLIVE_700 = RGBColor(0x44, 0x52, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"   # fallback robusto in PowerPoint per il serif (Fraunces non sempre installato)
SANS = "Calibri"    # fallback robusto per Inter
# Nota: se l'utente ha Fraunces/Inter installati, può sostituirli nel master.

# Dimensioni slide 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_solid_fill(shape, fill_color)
    if not line:
        set_no_line(shape)
    return shape


def add_text(slide, x, y, w, h, text, *,
             font=SANS, size=14, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # text può contenere \n per a-capo
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            # spaziatura tra lettere via XML
            from pptx.oxml.ns import qn
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb


def add_eyebrow(slide, x, y, w, text):
    return add_text(slide, x, y, w, Inches(0.3),
                    text.upper(), font=SANS, size=10, color=OLIVE_700,
                    bold=True, letter_spacing=300)


def add_brand_mark(slide, x, y, width=Inches(0.6)):
    return add_rect(slide, x, y, width, Emu(38100), OLIVE_500)  # ~3pt high


def add_slide_number(slide, idx, total):
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.3),
             f"{idx:02d} / {total:02d}",
             font=SANS, size=9, color=INK_MUTED, align=PP_ALIGN.RIGHT,
             letter_spacing=200)


def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_paragraph_run(p, text, *, font=SANS, size=14, color=INK,
                       bold=False, italic=False):
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


# ---------- SLIDE BUILDERS ----------

def slide_cover(prs, idx, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM_50)

    cx = SLIDE_W / 2

    # brand mark centrato
    mark_w = Inches(1)
    add_rect(slide, cx - mark_w / 2, Inches(2.4), mark_w, Emu(50800), OLIVE_500)

    # Titolo IMercati
    add_text(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(2),
             "IMercati", font=SERIF, size=120, color=INK, bold=False,
             align=PP_ALIGN.CENTER)

    # Tagline
    add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(1),
             "I mercati della provincia di Imperia. Insieme.",
             font=SERIF, size=24, color=INK_SOFT, italic=True,
             align=PP_ALIGN.CENTER)

    # Footer
    add_text(slide, Inches(0.8), Inches(6.9), Inches(4), Inches(0.4),
             "Maggio 2026", font=SANS, size=10, color=INK_MUTED,
             letter_spacing=200)
    add_text(slide, Inches(8.5), Inches(6.9), Inches(4), Inches(0.4),
             "Presentazione ai primi 10 operatori",
             font=SANS, size=10, color=INK_MUTED,
             align=PP_ALIGN.RIGHT, letter_spacing=200)

    set_notes(slide, "Buonasera a tutti. Vi ho chiamati perché siete dieci "
              "persone che mio padre conosce e di cui si fida. E perché "
              "secondo me siete i dieci che possono far partire una cosa che, "
              "da soli, nessuno di noi può fare. Vi prendo venti minuti, poi "
              "parliamo.")
    return slide


def slide_with_header(prs, idx, total, eyebrow, title_lines,
                      title_italic_parts=None, bg=CREAM_50):
    """Crea slide standard con header (eyebrow + titolo serif a due righe).
    title_lines: lista di stringhe (una per riga). Parti tra * * vengono
    rese in italico + olive_600.
    """
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg)

    x = Inches(0.9)
    add_eyebrow(slide, x, Inches(0.7), Inches(8), eyebrow)

    # Titolo: textbox unico, ogni riga un paragrafo
    tb = slide.shapes.add_textbox(x, Inches(1.1), Inches(11.5), Inches(2.2))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    for i, line in enumerate(title_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Parsing semplice: *text* → italic olive
        parts = []
        buf = ""
        in_em = False
        for ch in line:
            if ch == "*":
                if buf:
                    parts.append((buf, in_em))
                    buf = ""
                in_em = not in_em
            else:
                buf += ch
        if buf:
            parts.append((buf, in_em))
        for text, em in parts:
            add_paragraph_run(p, text, font=SERIF, size=54,
                              color=OLIVE_600 if em else INK,
                              italic=em, bold=False)

    add_slide_number(slide, idx, total)
    return slide


def fact_list(slide, items, y_start=Inches(3.6)):
    """Lista 'fatti' big-serif con trattino olive."""
    y = y_start
    for item in items:
        # trattino
        add_rect(slide, Inches(0.9), y + Inches(0.35), Inches(0.25),
                 Emu(12700), OLIVE_500)
        add_text(slide, Inches(1.3), y, Inches(11), Inches(0.7),
                 item, font=SERIF, size=26, color=INK_SOFT)
        y += Inches(0.9)


def bullet_list(slide, items, y_start=Inches(3.6), bullet_color=OLIVE_500,
                font_size=18, x_start=Inches(0.9), width=Inches(11.5)):
    """Lista bullet con pallino olive. items: lista di tuple (text, bold_parts)
    o solo stringhe."""
    y = y_start
    for item in items:
        text = item if isinstance(item, str) else item[0]
        # pallino
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_start + Inches(0.05),
                                      y + Inches(0.18), Inches(0.12),
                                      Inches(0.12))
        set_solid_fill(dot, bullet_color)
        set_no_line(dot)
        # testo con possibili **bold**
        tb = slide.shapes.add_textbox(x_start + Inches(0.45), y,
                                       width - Inches(0.45), Inches(0.7))
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # parsing **bold**
        parts = []
        buf = ""
        in_bold = False
        i = 0
        while i < len(text):
            if text[i:i+2] == "**":
                if buf:
                    parts.append((buf, in_bold))
                    buf = ""
                in_bold = not in_bold
                i += 2
            else:
                buf += text[i]
                i += 1
        if buf:
            parts.append((buf, in_bold))
        for txt, is_bold in parts:
            add_paragraph_run(p, txt, font=SANS, size=font_size,
                              color=INK if is_bold else INK_SOFT,
                              bold=is_bold)
        y += Inches(0.65)


def big_statement(slide, text, y_start=Inches(3.5), size=28,
                   color=INK, italic_marker=True):
    """Statement importante, con *...* per italic olive."""
    tb = slide.shapes.add_textbox(Inches(0.9), y_start, Inches(11.5),
                                   Inches(3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    for li, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        if not line.strip():
            add_paragraph_run(p, " ", font=SERIF, size=size, color=color)
            continue
        parts = []
        buf = ""
        in_em = False
        for ch in line:
            if ch == "*":
                if buf:
                    parts.append((buf, in_em))
                    buf = ""
                in_em = not in_em
            else:
                buf += ch
        if buf:
            parts.append((buf, in_em))
        for txt, em in parts:
            add_paragraph_run(p, txt, font=SERIF, size=size,
                              color=OLIVE_600 if em else color, italic=em)


# ---------- SINGLE SLIDES ----------

def build_slide_02(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Il punto di partenza",
                              ["Lo sappiamo", "tutti."])
    fact_list(slide, [
        "I clienti dei mercati calano ogni anno.",
        "I giovani comprano online, gli anziani escono di meno.",
        "Chi entra al mercato per caso è sempre meno.",
    ])
    set_notes(slide, "Non vi vendo niente. Vi sto dicendo quello che vedete "
              "voi tutti i sabati mattina. Non è colpa vostra, è il mondo che "
              "è cambiato. Però non vuol dire che dobbiamo subire e basta.")
    return slide


def build_slide_03(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Il vero problema",
                              ['Per i clienti', 'siete *"il banco".*'])
    big_statement(slide,
                  "Non Mario. Non Maria.\n"
                  "Non chi pesca il pesce alle quattro del mattino.\n"
                  "Non chi sceglie l'olio dei contadini di Pontedassio.\n"
                  "\n"
                  '*"Il banco della frutta".*\n'
                  "E al banco — si sostituisce facile.",
                  y_start=Inches(3.4), size=24)
    set_notes(slide, "Pausa. Lasciali pensare. È la slide più importante "
              "della presentazione. Il problema non è il calo. Il problema è "
              "che siete invisibili come persone. Per il cliente siete 'il "
              "banco della frutta'. Non Mario, non Maria. Il banco. E al "
              "banco si sostituisce facile.")
    return slide


def build_slide_04(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Il lavoro vero",
                              ["Costruire il vostro",
                               "*posizionamento.*"])
    big_statement(slide,
                  "Una parola difficile per una cosa semplice:\n"
                  "\n"
                  '*"Perché dovrei comprare da te,*\n'
                  '*e non dal banco accanto?"*\n'
                  "\n"
                  "Se la risposta è solo il prezzo,\n"
                  "avete già perso. C'è il supermercato.",
                  y_start=Inches(3.3), size=22)
    set_notes(slide, "Questa è la parola chiave di oggi: POSIZIONAMENTO. Non "
              "è marketing astratto. È la risposta a una domanda semplice: "
              "'perché un cliente nuovo dovrebbe scegliere te e non il banco "
              "accanto?'. Se la risposta è solo il prezzo, avete perso in "
              "partenza. C'è sempre il supermercato.")
    return slide


def build_slide_05(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Però...",
                              ["Avete tutto quello",
                               "che serve. Già."])
    bullet_list(slide, [
        "Un **mestiere** che fate da vent'anni — quello non si improvvisa.",
        "Una **faccia, una voce, una storia**. Tre cose che Amazon non avrà mai.",
        "Conoscete **i contadini, i pescatori, i fornitori** per nome.",
        "Sapete **scegliere il prodotto** meglio del cliente — perché lo fate da una vita.",
        "Lo date per scontato. **Non lo raccontate.** Questo è l'errore.",
    ], y_start=Inches(3.5), font_size=18)
    set_notes(slide, "Adesso ribaltiamo. Avete TUTTO quello che serve per "
              "costruire un posizionamento forte. Lo avete già. Solo che non "
              "lo raccontate. Lo date per scontato. Per voi è normale alzarvi "
              "alle quattro, per il cliente è una storia. Per voi è normale "
              "conoscere il contadino, per il cliente è la garanzia.")
    return slide


def build_slide_06_overview(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Come ci arriviamo",
                              ["Quattro tappe.",
                               "Voi al centro."])
    # Sottotitolo
    add_text(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.5),
             "Vi accompagno all'inizio. Dopo, camminate da soli.",
             font=SERIF, size=17, color=INK_SOFT, italic=True)

    # 4 step cards in griglia
    card_w = Inches(2.85)
    card_h = Inches(3)
    gap = Inches(0.25)
    start_x = Inches(0.9)
    y = Inches(3.85)

    steps = [
        ("1", "Mese 1 — Capirvi", "Chi siete davvero",
         "Form, una chiacchierata, mettiamo per iscritto la vostra storia e cosa vi rende diversi."),
        ("2", "Mese 1–2 — Metterlo a fuoco", "Il vostro profilo",
         "Costruiamo insieme la prima versione del messaggio: foto, parole, primi contenuti."),
        ("3", "Mese 2–4 — Provare", "Pubblicare e ascoltare",
         "Mettete fuori contenuti, vediamo cosa funziona, aggiustiamo. Faccio io, con voi."),
        ("4", "Mese 4+ — Autonomia", "Camminate da soli",
         "Avete il metodo, gli strumenti, le abitudini. Io resto disponibile, ma non vi servo più ogni giorno."),
    ]
    for i, (num, when, title, desc) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        # card bianca
        card = add_rect(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = CREAM_300
        # barra olive sinistra
        add_rect(slide, x, y, Inches(0.06), card_h, OLIVE_500)

        # numero grande
        add_text(slide, x + Inches(0.25), y + Inches(0.15),
                 card_w - Inches(0.3), Inches(0.9),
                 num, font=SERIF, size=42, color=OLIVE_500)
        # when
        add_text(slide, x + Inches(0.25), y + Inches(1.0),
                 card_w - Inches(0.3), Inches(0.4),
                 when.upper(), font=SANS, size=8, color=INK_MUTED,
                 bold=True, letter_spacing=200)
        # titolo
        add_text(slide, x + Inches(0.25), y + Inches(1.35),
                 card_w - Inches(0.3), Inches(0.8),
                 title, font=SERIF, size=18, color=INK, bold=True)
        # desc
        add_text(slide, x + Inches(0.25), y + Inches(2.0),
                 card_w - Inches(0.3), Inches(0.95),
                 desc, font=SANS, size=10, color=INK_SOFT)

    set_notes(slide, "Adesso vi mostro come ci arriviamo. Quattro tappe. "
              "Quattro mesi. Alla fine voi dovete essere AUTONOMI — non "
              "dipendenti da me. Io sono il facilitatore dei primi mesi, poi "
              "voi camminate. Questa è la promessa principale dell'incontro "
              "di oggi.")
    return slide


def build_step_zoom(prs, idx, total, step_num, step_when, eyebrow,
                     title_lines, lead, bullets, my_role, notes):
    """Slide zoom su un singolo step."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM_50)

    # Eyebrow
    add_eyebrow(slide, Inches(0.9), Inches(0.7), Inches(8), eyebrow)

    # Badge gigante a sinistra
    add_text(slide, Inches(0.9), Inches(1.5), Inches(3.5), Inches(2.5),
             step_num, font=SERIF, size=120, color=OLIVE_500)
    add_text(slide, Inches(0.9), Inches(3.6), Inches(3.5), Inches(0.4),
             step_when.upper(), font=SANS, size=9, color=INK_MUTED,
             bold=True, letter_spacing=300)

    # Titolo destra
    tb = slide.shapes.add_textbox(Inches(4.8), Inches(1.3), Inches(8),
                                   Inches(1.6))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    for i, line in enumerate(title_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        parts = []
        buf = ""
        in_em = False
        for ch in line:
            if ch == "*":
                if buf:
                    parts.append((buf, in_em))
                    buf = ""
                in_em = not in_em
            else:
                buf += ch
        if buf:
            parts.append((buf, in_em))
        for txt, em in parts:
            add_paragraph_run(p, txt, font=SERIF, size=34,
                              color=OLIVE_600 if em else INK, italic=em)

    # Lead line
    add_text(slide, Inches(4.8), Inches(3.0), Inches(8), Inches(0.5),
             lead, font=SERIF, size=15, color=INK_SOFT, italic=True)

    # Bullets
    bullet_list(slide, bullets, y_start=Inches(3.7), font_size=14,
                x_start=Inches(4.8), width=Inches(7.8))

    # Divisorio + my role
    role_y = Inches(6.0)
    add_rect(slide, Inches(4.8), role_y, Inches(7.8), Emu(6350), CREAM_300)
    add_text(slide, Inches(4.8), role_y + Inches(0.12), Inches(8),
             Inches(0.35), "COSA FACCIO IO",
             font=SANS, size=9, color=OLIVE_700, bold=True,
             letter_spacing=300)
    add_text(slide, Inches(4.8), role_y + Inches(0.5), Inches(7.8),
             Inches(0.8), my_role,
             font=SERIF, size=12, color=INK_MUTED, italic=True)

    add_slide_number(slide, idx, total)
    set_notes(slide, notes)
    return slide


def build_slide_11_assets(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Su cosa ci appoggiamo",
                              ["Gli asset della",
                               "*strategia.*"])

    card_w = Inches(3.85)
    card_h = Inches(3)
    gap = Inches(0.3)
    start_x = Inches(0.9)
    y = Inches(3.7)

    assets = [
        (True, "Asset principale", "Voi.",
         "La vostra faccia, la vostra storia, il vostro mestiere. "
         "Tutto il resto è al servizio di questo."),
        (False, "Strumento", "I vostri canali",
         "Facebook, Instagram, WhatsApp. Quelli che già usate. "
         "Li valorizziamo, non ne aggiungiamo di nuovi."),
        (False, "Strumento", "IMercati",
         "Il sito dove la gente vi cerca. Profilo, foto, prodotti, "
         "orari. È già online — lo guardate dopo."),
    ]
    for i, (primary, label, title, desc) in enumerate(assets):
        x = start_x + (card_w + gap) * i
        bg = OLIVE_50 if primary else CREAM_100
        card = add_rect(slide, x, y, card_w, card_h, bg)
        card.line.color.rgb = OLIVE_300 if primary else CREAM_300

        add_text(slide, x + Inches(0.4), y + Inches(0.35),
                 card_w - Inches(0.5), Inches(0.4),
                 label.upper(), font=SERIF, size=10, color=OLIVE_600,
                 bold=True, letter_spacing=200)
        add_text(slide, x + Inches(0.4), y + Inches(0.85),
                 card_w - Inches(0.5), Inches(0.7),
                 title, font=SERIF, size=24, color=INK, bold=True)
        add_text(slide, x + Inches(0.4), y + Inches(1.65),
                 card_w - Inches(0.5), Inches(1.3),
                 desc, font=SANS, size=12, color=INK_SOFT)

    set_notes(slide, "Tutto questo percorso si appoggia su tre asset. Il "
              "vostro brand è il principale — siete voi, la vostra storia, "
              "la vostra faccia. Poi i social personali, che già qualcuno di "
              "voi usa. E poi IMercati: il SITO è lo strumento che mette "
              "ordine, ma è uno strumento, non IL progetto. Voglio che lo "
              "abbiate chiaro: il sito serve a voi, non viceversa.")
    return slide


def build_slide_12_patto(prs, idx, total):
    slide = slide_with_header(prs, idx, total,
                              "Il patto",
                              ["Tra noi."])

    # Due colonne
    col_w = Inches(5.7)
    gap = Inches(0.3)
    x1 = Inches(0.9)
    x2 = x1 + col_w + gap
    y = Inches(3.4)

    # Colonna 1
    add_text(slide, x1, y, col_w, Inches(0.6), "Io mi impegno a",
             font=SERIF, size=22, color=INK, bold=True)
    add_rect(slide, x1, y + Inches(0.65), col_w, Emu(25400), OLIVE_500)
    bullet_list(slide, [
        "Accompagnarvi **per quattro mesi**, fino all'autonomia.",
        "**Non chiedervi soldi**, mai, voi dieci.",
        "Darvi **metodo, strumenti, abitudini** — non solo contenuti.",
        "Restare disponibile **anche dopo**, su decisioni grandi.",
    ], y_start=y + Inches(0.95), font_size=14, x_start=x1, width=col_w)

    # Colonna 2
    add_text(slide, x2, y, col_w, Inches(0.6), "Voi vi impegnate a",
             font=SERIF, size=22, color=INK, bold=True)
    add_rect(slide, x2, y + Inches(0.65), col_w, Emu(25400), OLIVE_500)
    bullet_list(slide, [
        "Compilare il form, **oggi o entro una settimana**.",
        "**Partecipare attivamente** ai quattro step. Non passivamente.",
        "Provare a **pubblicare con costanza** nei mesi 2–4.",
        "Dirmi **la verità** quando vi chiedo cosa ne pensate.",
    ], y_start=y + Inches(0.95), font_size=14, x_start=x2, width=col_w)

    set_notes(slide, "Il patto. Non c'è da firmare niente di legale. È una "
              "stretta di mano. Funziona se ci credete davvero, altrimenti "
              "meglio dirmelo adesso. La cosa importante: io vi accompagno "
              "fino al quarto mese, poi siete voi a guidare.")
    return slide


def build_slide_13_cta(prs, idx, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, CREAM_50)

    # Eyebrow centrato
    add_eyebrow(slide, Inches(0.5), Inches(1.4), Inches(12.3),
                "Step 1 inizia adesso")
    # Sposto eyebrow al centro
    # (textbox è già allineato a left default; lo rifaccio centered)
    # Per semplicità lo ridisegno:
    # rimozione: skip — usiamo textbox custom
    # Titolo centrato
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12.3),
                                   Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_paragraph_run(p1, "Il form.", font=SERIF, size=66, color=INK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    add_paragraph_run(p2, "Dieci minuti.", font=SERIF, size=66,
                      color=OLIVE_600, italic=True)

    add_text(slide, Inches(2), Inches(4.4), Inches(9.3), Inches(0.8),
             "Il primo gesto concreto del percorso. Si riempie qui, insieme. "
             "Poi domande, dubbi, tutto.",
             font=SANS, size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)

    # 3 cta cards
    card_w = Inches(3.6)
    card_h = Inches(1.6)
    gap = Inches(0.3)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(5.4)

    cards = [
        ("Form cartaceo", "10 copie sul tavolo, pronte"),
        ("Versione online", "QR code (da inserire)"),
        ("Per dopo", "[Telefono]\n[Email]"),
    ]
    for i, (label, value) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        card = add_rect(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = CREAM_300
        add_text(slide, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.6),
                 Inches(0.5), label,
                 font=SERIF, size=16, color=INK, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), y + Inches(0.75), card_w - Inches(0.6),
                 Inches(0.8), value,
                 font=SANS, size=11, color=INK_MUTED,
                 align=PP_ALIGN.CENTER)

    add_slide_number(slide, idx, total)
    set_notes(slide, "E adesso si comincia. Step uno è esattamente questo: "
              "il form. Se siete con me, lo riempiamo qui, insieme. Dieci "
              "minuti. È il primo gesto concreto di tutto il percorso. Dopo, "
              "parliamo, domande, dubbi, tutto.")
    return slide


# ---------- MAIN ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 13

    slide_cover(prs, 1, total)
    build_slide_02(prs, 2, total)
    build_slide_03(prs, 3, total)
    build_slide_04(prs, 4, total)
    build_slide_05(prs, 5, total)
    build_slide_06_overview(prs, 6, total)

    # Step zooms
    build_step_zoom(
        prs, 7, total, "01", "Capirvi", "Step 1 — Mese 1",
        ["Mettere per iscritto",
         "quello che date per *scontato.*"],
        "Voi sapete tutto. Ma è 'implicito'. Lo facciamo diventare parole.",
        [
            "Compilate il **form di oggi** — 34 domande, dieci minuti.",
            "Una **chiacchierata di 30 minuti**, di persona o al telefono, per approfondire.",
            "Tiriamo fuori **tre cose**: chi siete, cosa vi rende diversi, la frase che vi descrive in dieci parole.",
        ],
        "Vi ascolto, scrivo, vi rimando una scheda chiara. Quella scheda è la base di tutto il resto.",
        "Step uno. Capire chi siete. Il form di oggi è esattamente questo — non "
        "è una registrazione burocratica. È il mio modo di farvi DIRE per "
        "iscritto cose che date per scontate. È la materia prima per tutto il "
        "resto. Senza, non si fa niente."
    )

    build_step_zoom(
        prs, 8, total, "02", "Metterlo a fuoco", "Step 2 — Mese 1–2",
        ["Il vostro profilo",
         "diventa *visibile.*"],
        "Costruiamo insieme la prima versione del messaggio. Imperfetta, ma vostra.",
        [
            "Una **sessione foto al banco** — vi raccontate al cliente, io scatto.",
            "Scriviamo insieme **il vostro testo**: chi siete, in una pagina.",
            "Pubblichiamo il profilo su **IMercati** e prepariamo **3 contenuti pronti** per i vostri social.",
        ],
        "Scrivo, scatto, monto. Vi mostro come l'ho fatto, così la seconda volta lo facciamo insieme.",
        "Step due. Dal sapere chi siete a dirlo bene. Qui costruiamo la prima "
        "versione di tutto: il vostro profilo sul sito, le prime tre foto "
        "buone, la prima volta che provate a scrivere qualcosa sui vostri "
        "canali. Non aspettiamo la perfezione — partiamo. Si aggiusta strada "
        "facendo."
    )

    build_step_zoom(
        prs, 9, total, "03", "Provare", "Step 3 — Mese 2–4",
        ["Pubblicare,",
         "ascoltare, *aggiustare.*"],
        "Si esce in strada. Si guarda cosa funziona. Si cambia se serve.",
        [
            "Un **contenuto a settimana** minimo — facciamo insieme, voi sempre più protagonisti.",
            "Una **chiamata ogni due settimane** di 30 minuti per decidere cosa funziona.",
            "Imparate a **riconoscere** quando un contenuto 'tira' e quando no — è esperienza, non magia.",
        ],
        "Vi mostro i dati, vi do consigli, ma le decisioni sono vostre. Inizia il passaggio di mano.",
        "Step tre. Adesso si esce in strada. Pubblicate i contenuti, ascoltate "
        "cosa dicono i clienti, vediamo i numeri. Questa è la fase in cui "
        "passate dal SUBIRE al GUIDARE. Io vi affianco ma cominciate a "
        "prendere decisioni da soli. Imparate a riconoscere cosa funziona."
    )

    build_step_zoom(
        prs, 10, total, "04", "Autonomia", "Step 4 — Mese 4 in poi",
        ["Camminate *da soli.*",
         "Questo è il vero obiettivo."],
        "Avete il metodo, gli strumenti, le abitudini. Io resto, ma non vi servo più ogni giorno.",
        [
            "Sapete **fare e pubblicare contenuti** da soli, anche solo col telefono.",
            "Avete uno **strumento (in arrivo) per generare post** partendo da una foto e due frasi.",
            "Il vostro **profilo IMercati** resta vivo, lo aggiornate voi quando serve.",
            "Io resto disponibile per **consigli, controlli, decisioni grandi**. Non per la routine.",
        ],
        "Mi tolgo di mezzo. Resto allenatore — non più assistente.",
        "Step quattro. Questo è il vero obiettivo di tutta la cosa. Dopo "
        "quattro mesi, voi camminate. Non perché vi ho mollato — perché siete "
        "VOI a saperlo fare. È come imparare a guidare: all'inizio uno in "
        "macchina con voi, poi guidate da soli. Io resto disponibile, ma non "
        "sono più necessario tutti i giorni."
    )

    build_slide_11_assets(prs, 11, total)
    build_slide_12_patto(prs, 12, total)
    build_slide_13_cta(prs, 13, total)

    out = Path(__file__).parent / "presentazione.pptx"
    prs.save(str(out))
    print(f"OK -> {out}")


if __name__ == "__main__":
    main()
